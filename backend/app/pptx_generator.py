import os
import logging
import re
from typing import List, Set, Tuple, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap
from dotenv import load_dotenv
from .models import SlideContent
import requests
from io import BytesIO
from PIL import Image

# Load environment variables from backend/.env
backend_dir = Path(__file__).resolve().parent.parent
dotenv_path = backend_dir / '.env'
load_dotenv(dotenv_path=dotenv_path)

logger = logging.getLogger(__name__)


def extract_image_colors(image_stream: BytesIO) -> List[Tuple[int, int, int]]:
    """
    Extract key colors from image at center and four corners with weighted averaging.
    Returns list of RGB tuples for sampling points.
    """
    try:
        image_stream.seek(0)
        img = Image.open(image_stream)
        img = img.convert('RGB')
        width, height = img.size
        
        # Define sampling points: center (weight=3), four corners (weight=1 each)
        sample_points = [
            (width // 2, height // 2, 3.0),      # Center (highest weight)
            (10, 10, 1.0),                        # Top-left
            (width - 10, 10, 1.0),                # Top-right
            (10, height - 10, 1.0),               # Bottom-left
            (width - 10, height - 10, 1.0),       # Bottom-right
            (width // 2, 10, 1.5),                # Top-center (medium weight)
            (width // 2, height - 10, 1.5),       # Bottom-center (medium weight)
        ]
        
        colors_with_weights = []
        for x, y, weight in sample_points:
            # Sample a small region around each point for stability
            x = max(5, min(x, width - 5))
            y = max(5, min(y, height - 5))
            
            # Get average color of 5x5 region
            region_colors = []
            for dx in range(-2, 3):
                for dy in range(-2, 3):
                    px, py = max(0, min(x + dx, width - 1)), max(0, min(y + dy, height - 1))
                    region_colors.append(img.getpixel((px, py)))
            
            avg_r = sum(c[0] for c in region_colors) // len(region_colors)
            avg_g = sum(c[1] for c in region_colors) // len(region_colors)
            avg_b = sum(c[2] for c in region_colors) // len(region_colors)
            colors_with_weights.append(((avg_r, avg_g, avg_b), weight))
        
        image_stream.seek(0)  # Reset stream for later use
        return colors_with_weights
    except Exception as e:
        logger.error(f"Failed to extract image colors: {e}")
        return [((128, 128, 128), 1.0)]  # Default gray


def calculate_weighted_background_color(colors_with_weights: List[Tuple[Tuple[int, int, int], float]]) -> RGBColor:
    """
    Calculate weighted average background color from sampled points.
    """
    if not colors_with_weights:
        return RGBColor(128, 128, 128)
    
    total_weight = sum(w for _, w in colors_with_weights)
    weighted_r = sum(c[0] * w for c, w in colors_with_weights) / total_weight
    weighted_g = sum(c[1] * w for c, w in colors_with_weights) / total_weight
    weighted_b = sum(c[2] * w for c, w in colors_with_weights) / total_weight
    
    return RGBColor(int(weighted_r), int(weighted_g), int(weighted_b))


def calculate_luminance(rgb_color: RGBColor) -> float:
    """
    Calculate relative luminance of a color (0.0 = darkest, 1.0 = brightest).
    Uses WCAG formula for perceptual luminance.
    """
    r, g, b = rgb_color[0] / 255.0, rgb_color[1] / 255.0, rgb_color[2] / 255.0
    
    # Apply gamma correction
    r = r / 12.92 if r <= 0.03928 else ((r + 0.055) / 1.055) ** 2.4
    g = g / 12.92 if g <= 0.03928 else ((g + 0.055) / 1.055) ** 2.4
    b = b / 12.92 if b <= 0.03928 else ((b + 0.055) / 1.055) ** 2.4
    
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def get_contrasting_color_with_shadow(background_color: RGBColor) -> Tuple[RGBColor, RGBColor]:
    """
    Returns text color and shadow color based on background luminance.
    Dark backgrounds get white text with dark shadow, light backgrounds get dark text with light shadow.
    """
    luminance = calculate_luminance(background_color)
    
    if luminance < 0.5:
        # Dark background: white text with semi-transparent black shadow
        text_color = RGBColor(255, 255, 255)
        shadow_color = RGBColor(0, 0, 0)
    else:
        # Light background: dark text with semi-transparent white shadow
        text_color = RGBColor(30, 30, 30)
        shadow_color = RGBColor(255, 255, 255)
    
    return text_color, shadow_color


def get_contrasting_color(background_color: RGBColor) -> RGBColor:
    """
    Returns white or black text color based on background luminance.
    Dark backgrounds get white text, light backgrounds get black text.
    """
    luminance = calculate_luminance(background_color)
    # If background is dark (luminance < 0.7), use white text
    # If background is light (luminance >= 0.7), use black text
    return RGBColor(255, 255, 255) if luminance < 0.7 else RGBColor(0, 0, 0)


class PPTXGenerator:
    """
    Generator for PowerPoint presentations using python-pptx.
    Handles template selection, slide creation, and rich content rendering.
    """

    def __init__(self):
        # Resolve absolute paths relative to this file: backend/app/pptx_generator.py -> backend/
        base_dir = Path(__file__).resolve().parent.parent
        
        default_output = str(base_dir / "output")
        default_templates = str(base_dir / "templates")

        self.output_dir = os.getenv("OUTPUT_DIR", default_output)
        self.templates_dir = os.getenv("TEMPLATES_DIR", default_templates)

        os.makedirs(self.output_dir, exist_ok=True)

        # Template color schemes
        self.color_schemes = {
            "corporate": {
                "primary": RGBColor(0, 51, 102),      # Dark blue
                "secondary": RGBColor(0, 102, 204),    # Medium blue
                "accent": RGBColor(255, 153, 0)        # Orange
            },
            "academic": {
                "primary": RGBColor(51, 51, 51),       # Dark gray
                "secondary": RGBColor(102, 102, 102),  # Medium gray
                "accent": RGBColor(153, 0, 0)          # Maroon
            },
            "startup": {
                "primary": RGBColor(102, 45, 145),     # Purple
                "secondary": RGBColor(153, 102, 255),  # Light purple
                "accent": RGBColor(255, 195, 0)        # Gold
            },
            "minimal": {
                "primary": RGBColor(0, 0, 0),          # Black
                "secondary": RGBColor(128, 128, 128),  # Gray
                "accent": RGBColor(255, 255, 255)      # White
            }
        }

    def _extract_keywords(self, content: SlideContent) -> Set[str]:
        """
        Extract relevant keywords from slide content for image search.
        Returns a set of meaningful keywords.
        """
        keywords = set()
        
        # Extract from title
        if content.title:
            # Remove common stop words and split
            title_words = re.findall(r'\b\w+\b', content.title.lower())
            stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'is', 'are', 'was', 'were', 'be', 'been', 'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would', 'could', 'should', 'may', 'might', 'must', 'can', 'this', 'that', 'these', 'those'}
            title_keywords = [word for word in title_words if len(word) > 2 and word not in stop_words]
            keywords.update(title_keywords[:5])  # Limit to top 5 keywords from title
        
        # Extract from content/paragraph
        text_content = ""
        if content.content:
            text_content += " ".join(content.content)
        if content.paragraph:
            text_content += " " + content.paragraph
            
        if text_content:
            content_words = re.findall(r'\b\w+\b', text_content.lower())
            content_keywords = [word for word in content_words if len(word) > 3 and word not in stop_words]
            # Count frequency and get most common
            from collections import Counter
            word_counts = Counter(content_keywords)
            top_keywords = [word for word, count in word_counts.most_common(3) if count > 1]
            keywords.update(top_keywords)
        
        # Extract from image description if available
        if content.image_description:
            desc_words = re.findall(r'\b\w+\b', content.image_description.lower())
            desc_keywords = [word for word in desc_words if len(word) > 2 and word not in stop_words]
            keywords.update(desc_keywords[:3])
            
        return keywords

    def _is_image_relevant(self, keywords: Set[str], image_description: str = None, image_url: str = None) -> bool:
        """
        Check if an image is relevant to the slide content based on keywords.
        """
        if not keywords:
            return False
            
        # If we have image description, check keyword overlap
        if image_description:
            desc_text = image_description.lower()
            matching_keywords = [kw for kw in keywords if kw in desc_text]
            return len(matching_keywords) >= 1  # At least one keyword match
        
        # If we have image URL, try to extract keywords from URL
        if image_url:
            url_lower = image_url.lower()
            matching_keywords = [kw for kw in keywords if kw in url_lower]
            return len(matching_keywords) >= 1
            
        return False

    def _has_background_image(self, slide) -> bool:
        """
        Check if the slide already has a background image.
        """
        try:
            # Check if background has an image fill
            bg = slide.background
            if hasattr(bg, 'fill') and bg.fill.type == 2:  # Picture fill type
                return True
            # Check for shapes that might be background images (large images covering most of the slide)
            for shape in slide.shapes:
                if hasattr(shape, 'image') and shape.width >= Inches(9) and shape.height >= Inches(6):
                    return True
        except Exception:
            pass
        return False

    async def create_presentation(
        self,
        deck_id: str,
        slides: List[SlideContent],
        template: str,
        title: str,
        design_config: dict = None
    ) -> str:
        logger.info(f"Creating presentation with {len(slides)} slides")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Determine colors
        colors = self.color_schemes.get("corporate").copy() 
        
        if design_config and "colors" in design_config:
            try:
                agent_colors = design_config["colors"]
                new_colors = {}
                for key, rgb_list in agent_colors.items():
                    if isinstance(rgb_list, list) and len(rgb_list) == 3:
                        new_colors[key] = RGBColor(int(rgb_list[0]), int(rgb_list[1]), int(rgb_list[2]))
                colors.update(new_colors)
            except Exception as e:
                logger.error(f"Failed to apply design colors: {e}")

        # Add slides
        for i, slide_content in enumerate(slides):
            layout_idx = 1 # Default
            if slide_content.layout:
                layout_idx = slide_content.layout.layout_idx
            elif i == 0 or slide_content.slideType == "title":
                layout_idx = 0
            
            try:
                layout = prs.slide_layouts[layout_idx]
            except:
                layout = prs.slide_layouts[1] 
                
            slide = prs.slides.add_slide(layout)
            
            # Set background (image or color) and get detected background color
            detected_bg_color = None
            if slide_content.background_image_url:
                detected_bg_color = self._set_slide_background_image(slide, slide_content.background_image_url)
            else:
                self._set_slide_background(slide, colors)
            
            # Use detected background color for text contrast if available
            slide_colors = colors.copy()
            if detected_bg_color:
                slide_colors["background"] = detected_bg_color
            
            self._fill_slide_content(slide, slide_content, slide_colors, layout_idx)

        file_path = os.path.join(self.output_dir, f"{deck_id}.pptx")
        prs.save(file_path)
        logger.info(f"Presentation saved: {file_path}")
        return file_path

    def _set_slide_background(self, slide, colors: dict):
        if "background" in colors:
            try:
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = colors["background"]
            except Exception:
                pass

    def _fill_slide_content(self, slide, content: SlideContent, colors: dict, layout_idx: int):
        # Determine text color based on background
        bg_color = colors.get("background", colors.get("primary", RGBColor(255, 255, 255)))
        auto_text_color = get_contrasting_color(bg_color)
        
        # 1. Set Title
        if slide.shapes.title:
            slide.shapes.title.text = content.title
            for paragraph in slide.shapes.title.text_frame.paragraphs:
                paragraph.font.size = Pt(40) if layout_idx == 0 else Pt(32)
                paragraph.font.color.rgb = auto_text_color
                paragraph.font.bold = True
                
        # 2. Main Content Logic
        
        # A. TABLE SLIDE
        if content.slideType == "table" and content.table:
            self._add_table(slide, content.table, colors)
            return

        # B. IMAGE SLIDE - dedicated image slides
        if content.slideType == "image" or content.image_description:
            # Check if slide already has background image - don't add another image
            has_background = self._has_background_image(slide)
            
            if not has_background and content.image_url:
                # Extract keywords and check relevance
                keywords = self._extract_keywords(content)
                if self._is_image_relevant(keywords, content.image_description, content.image_url):
                    self._add_image_from_url(slide, content.image_url)
                    logger.info(f"Added relevant image for slide: {content.title}")
                else:
                    logger.info(f"Skipped irrelevant image for slide: {content.title}")
                    self._add_image_placeholder(slide, content.image_description or "Image not relevant to content", colors)
            elif not has_background:
                # Fallback to placeholder
                self._add_image_placeholder(slide, content.image_description, colors)
            
            # Add text if exists
            if content.content and len(slide.placeholders) > 1:
                 # Check if placeholder 1 is text
                 try: 
                    self._populate_text_frame(slide.placeholders[1].text_frame, content.content, colors) 
                 except: pass
            return

        # C. NARRATIVE / PARAGRAPH
        if content.paragraph:
            # Check if slide already has background image
            has_background = self._has_background_image(slide)
            
            # If image available and no background image, check relevance before adding
            should_add_image = False
            if content.image_url and not has_background:
                keywords = self._extract_keywords(content)
                if self._is_image_relevant(keywords, content.image_description, content.image_url):
                    should_add_image = True
                    logger.info(f"Added relevant image for paragraph slide: {content.title}")
                else:
                    logger.info(f"Skipped irrelevant image for paragraph slide: {content.title}")
            
            if should_add_image:
                self._add_text_with_image_layout(slide, content, colors)
            else:
                # Use main body placeholder if available
                if len(slide.placeholders) > 1:
                    body = slide.placeholders[1]
                    if hasattr(body, "text_frame"):
                        self._populate_paragraph(body.text_frame, content.paragraph, colors)
                else:
                    # Fallback: Create custom textbox for Layout 5 (Title Only) or similar
                    left = Inches(1.0)
                    top = Inches(2.0)
                    width = Inches(8.0)
                    height = Inches(5.0)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    self._populate_paragraph(txBox.text_frame, content.paragraph, colors)
                    txBox.text_frame.word_wrap = True
            
            return

        # D. STANDARD LISTS (Fallback)
        if layout_idx == 0: # Title Slide
            if len(slide.placeholders) > 1 and content.content:
                subtitle = slide.placeholders[1]
                subtitle.text = content.content[0]
                for p in subtitle.text_frame.paragraphs:
                    p.font.size = Pt(20)
                    p.font.color.rgb = auto_text_color
                    
        elif layout_idx == 1: # Title + Content
            has_background = self._has_background_image(slide)
            
            if content.image_url and not has_background:
                # Check relevance before adding image
                keywords = self._extract_keywords(content)
                if self._is_image_relevant(keywords, content.image_description, content.image_url):
                    # Layout: text on left, image on right
                    self._add_list_with_image_layout(slide, content, colors)
                    logger.info(f"Added relevant image for list slide: {content.title}")
                else:
                    logger.info(f"Skipped irrelevant image for list slide: {content.title}")
                    # Fall back to text only
                    if len(slide.placeholders) > 1:
                        self._populate_text_frame(slide.placeholders[1].text_frame, content.content, colors)
            elif len(slide.placeholders) > 1:
                self._populate_text_frame(slide.placeholders[1].text_frame, content.content, colors)
                    
        elif layout_idx == 3: # Two Content - Optimized for long content
            # Check if we have long content that needs special handling
            has_long_content = False
            if content.content:
                has_long_content = any(len(point) > 20 for point in content.content)
            elif content.paragraph and len(content.paragraph) > 50:
                has_long_content = True
            
            if has_long_content:
                # For long content, use two-column layout with smaller font
                if content.content:
                    # Split bullet points into two columns
                    left_points = content.content[:len(content.content)//2]
                    right_points = content.content[len(content.content)//2:]
                    
                    if len(slide.placeholders) > 1:
                        self._populate_text_frame(slide.placeholders[1].text_frame, left_points, colors, force_small_font=True)
                    if len(slide.placeholders) > 2:
                        self._populate_text_frame(slide.placeholders[2].text_frame, right_points, colors, force_small_font=True)
                elif content.paragraph:
                    # Split paragraph into two columns
                    words = content.paragraph.split()
                    mid_point = len(words) // 2
                    left_text = " ".join(words[:mid_point])
                    right_text = " ".join(words[mid_point:])
                    
                    if len(slide.placeholders) > 1:
                        self._populate_paragraph(slide.placeholders[1].text_frame, left_text, colors, force_small_font=True)
                    if len(slide.placeholders) > 2:
                        self._populate_paragraph(slide.placeholders[2].text_frame, right_text, colors, force_small_font=True)
            else:
                # Standard two content layout
                left_points = content.content[:len(content.content)//2]
                right_points = content.content[len(content.content)//2:]
                
                if len(slide.placeholders) > 1:
                    self._populate_text_frame(slide.placeholders[1].text_frame, left_points, colors)
                if len(slide.placeholders) > 2:
                    self._populate_text_frame(slide.placeholders[2].text_frame, right_points, colors)
                
        else: # Generic Fallback
            if len(slide.placeholders) > 1:
                try: 
                    self._populate_text_frame(slide.placeholders[1].text_frame, content.content, colors)
                except AttributeError:
                    pass

    def _download_image(self, url: str) -> BytesIO:
        """Download image from URL and return as BytesIO object."""
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            return BytesIO(response.content)
        except Exception as e:
            logger.error(f"Failed to download image from {url}: {e}")
            return None
    
    def _set_slide_background_image(self, slide, image_url: str) -> Optional[RGBColor]:
        """
        Set slide background to an image from URL.
        Returns the weighted average background color for text contrast.
        """
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                # Extract colors for smart text contrast
                colors_with_weights = extract_image_colors(image_stream)
                weighted_bg_color = calculate_weighted_background_color(colors_with_weights)
                
                left = Inches(0)
                top = Inches(0)
                width = Inches(10)
                height = Inches(7.5)
                
                image_stream.seek(0)  # Reset stream position
                pic = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
                
                # Send to back (move to first position in shape collection)
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
                
                logger.info(f"Added background image from {image_url}, detected bg color: RGB({weighted_bg_color[0]}, {weighted_bg_color[1]}, {weighted_bg_color[2]})")
                return weighted_bg_color
        except Exception as e:
            logger.error(f"Failed to set background image: {e}")
        return None
    
    def _add_image_from_url(self, slide, image_url: str):
        """Add a large centered image from URL to the slide."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                # Center the image
                left = Inches(2.0)
                top = Inches(2.0)
                width = Inches(6.0)
                
                slide.shapes.add_picture(image_stream, left, top, width=width)
                logger.info(f"Added large image from {image_url}")
            else:
                # Fallback to placeholder
                self._add_image_placeholder(slide, "Image unavailable", {"secondary": RGBColor(200,200,200)})
        except Exception as e:
            logger.error(f"Failed to add image from URL: {e}")
            self._add_image_placeholder(slide, f"Image error: {str(e)[:50]}", {"secondary": RGBColor(200,200,200)})
    
    def _add_small_image(self, slide, image_url: str):
        """Add a smaller image as visual enhancement (right side of slide)."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                # Position on the right side
                left = Inches(6.5)
                top = Inches(2.5)
                width = Inches(3.0)
                
                slide.shapes.add_picture(image_stream, left, top, width=width)
                logger.info(f"Added small image from {image_url}")
        except Exception as e:
            logger.error(f"Failed to add small image: {e}")
    
    def _add_table(self, slide, table_data, colors):
        """Add a table to the slide."""
        rows = table_data.rows
        headers = table_data.headers
        if not rows or not headers: 
            return

        # Dimensions
        rows_count = len(rows) + 1
        cols_count = len(headers)
        
        # Centered table
        width = Inches(8.0)
        height = Inches(0.5 * rows_count + 0.5)
        left = Inches(1.0)
        top = Inches(2.0)

        shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
        table = shape.table

        # Style Headers
        for col, header_text in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = str(header_text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors.get("primary", RGBColor(0,0,0))
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = RGBColor(255,255,255)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER

        # Style Rows
        bg_color = colors.get("background", RGBColor(255, 255, 255))
        auto_text_color = get_contrasting_color(bg_color)
        
        for r, row_data in enumerate(rows):
            for c, cell_data in enumerate(row_data):
                if c < cols_count:
                    cell = table.cell(r+1, c)
                    cell.text = str(cell_data)
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(14)
                        p.font.color.rgb = auto_text_color
                        p.alignment = PP_ALIGN.LEFT

    def _download_image(self, url: str) -> BytesIO:
        """Download image from URL and return as BytesIO object."""
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            return BytesIO(response.content)
        except Exception as e:
            logger.error(f"Failed to download image from {url}: {e}")
            return None
    
    def _set_slide_background_image(self, slide, image_url: str):
        """Set slide background to an image from URL."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                # python-pptx doesn't directly support background images easily
                # Workaround: Add image as a shape covering the entire slide and send to back
                left = Inches(0)
                top = Inches(0)
                width = Inches(10)
                height = Inches(7.5)
                
                pic = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
                
                # Send to back (move to first position in shape collection)
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)  # Insert after background
                
                logger.info(f"Added background image from {image_url}")
        except Exception as e:
            logger.error(f"Failed to set background image: {e}")
    
    def _add_image_from_url(self, slide, image_url: str):
        """Add a large centered image from URL to the slide."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                # Center the image
                left = Inches(2.0)
                top = Inches(2.0)
                width = Inches(6.0)
                
                slide.shapes.add_picture(image_stream, left, top, width=width)
                logger.info(f"Added large image from {image_url}")
            else:
                # Fallback to placeholder
                self._add_image_placeholder(slide, "Image unavailable", {"secondary": RGBColor(200,200,200)})
        except Exception as e:
            logger.error(f"Failed to add image from URL: {e}")
            self._add_image_placeholder(slide, f"Image error: {str(e)[:50]}", {"secondary": RGBColor(200,200,200)})
    
    def _add_small_image(self, slide, image_url: str):
        """Add a smaller image as visual enhancement (right side of slide)."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                # Position on the right side
                left = Inches(6.5)
                top = Inches(2.5)
                width = Inches(3.0)
                
                slide.shapes.add_picture(image_stream, left, top, width=width)
                logger.info(f"Added small image from {image_url}")
        except Exception as e:
            logger.error(f"Failed to add small image: {e}")

    def _add_image_placeholder(self, slide, description, colors):
        """Add a placeholder shape for an image."""
        # Add a rounded rectangle
        left = Inches(1.5)
        top = Inches(2.5)
        width = Inches(7.0)
        height = Inches(4.0)
        
        # Check if we can put it in a specific placeholder?
        # For now, custom shape
        shape = slide.shapes.add_shape(
            1, # msoShapeRectangle
            left, top, width, height
        )
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = colors.get("secondary", RGBColor(200,200,200)) # Light gray-ish
        
        # Add text
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"[IMAGE PLACEHOLDER]\n{description}"
        p.alignment = PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(255,255,255)

    def _populate_paragraph(self, text_frame, text, colors, force_small_font=False):
        """Render a single narrative paragraph."""
        bg_color = colors.get("background", colors.get("primary", RGBColor(255, 255, 255)))
        auto_text_color = get_contrasting_color(bg_color)
        
        text_frame.clear()
        p = text_frame.add_paragraph()
        p.text = text
        # Use smaller font for long content or when forced
        font_size = Pt(18) if force_small_font else Pt(24)
        p.font.size = font_size
        p.font.color.rgb = auto_text_color
        p.alignment = PP_ALIGN.JUSTIFY
        p.space_after = Pt(10)

    def _populate_text_frame(self, text_frame, points, colors, force_small_font=False):
        bg_color = colors.get("background", colors.get("primary", RGBColor(255, 255, 255)))
        auto_text_color = get_contrasting_color(bg_color)
        luminance = calculate_luminance(bg_color)
        
        text_frame.clear()
        count = len(points)
        # Use smaller font for long content or when forced
        use_small_font = force_small_font or luminance < 0.5
        
        if count <= 3:
            base_size = 18 if use_small_font else 22
            spacing = 14
        elif count <= 5:
            base_size = 16 if use_small_font else 18
            spacing = 12
        else:
            base_size = 14 if use_small_font else 16
            spacing = 10

        for point in points:
            p = text_frame.add_paragraph()
            # "Key: Value" bolding logic
            if ":" in point and len(point.split(":")[0]) < 50:
                parts = point.split(":", 1)
                run_key = p.add_run()
                run_key.text = parts[0] + ":"
                run_key.font.bold = True
                run_key.font.size = Pt(base_size)
                run_key.font.color.rgb = auto_text_color
                
                run_val = p.add_run()
                run_val.text = parts[1]
                run_val.font.size = Pt(base_size)
                run_val.font.color.rgb = auto_text_color
            else:
                p.text = point
                p.font.size = Pt(base_size)
                p.font.color.rgb = auto_text_color
                if len(point) < 40 and count <= 4:
                    p.font.bold = True
            
            p.space_after = Pt(spacing)
            p.space_before = Pt(spacing/2)
    
    def _download_image(self, url: str) -> BytesIO:
        """Download image from URL and return as BytesIO object."""
        try:
            response = requests.get(url, timeout=10)
            response.raise_for_status()
            return BytesIO(response.content)
        except Exception as e:
            logger.error(f"Failed to download image from {url}: {e}")
            return None
    
    def _set_slide_background_image(self, slide, image_url: str):
        """Set slide background to an image from URL."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                left = Inches(0)
                top = Inches(0)
                width = Inches(10)
                height = Inches(7.5)
                
                pic = slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
                
                # Send to back
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.insert(2, pic._element)
                
                logger.info(f"Added background image from {image_url}")
        except Exception as e:
            logger.error(f"Failed to set background image: {e}")
    
    def _add_image_from_url(self, slide, image_url: str):
        """Add a large centered image from URL to the slide."""
        try:
            image_stream = self._download_image(image_url)
            if image_stream:
                left = Inches(2.0)
                top = Inches(2.0)
                width = Inches(6.0)
                
                slide.shapes.add_picture(image_stream, left, top, width=width)
                logger.info(f"Added large image from {image_url}")
            else:
                self._add_image_placeholder(slide, "Image unavailable", {"secondary": RGBColor(200,200,200)})
        except Exception as e:
            logger.error(f"Failed to add image from URL: {e}")
            self._add_image_placeholder(slide, f"Image error: {str(e)[:50]}", {"secondary": RGBColor(200,200,200)})
    
    def _add_text_with_image_layout(self, slide, content, colors):
        """Layout: paragraph text on left, image on right."""
        try:
            # Add text on left side
            left = Inches(0.8)
            top = Inches(2.0)
            width = Inches(5.0)
            height = Inches(5.0)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            self._populate_paragraph(text_box.text_frame, content.paragraph, colors)
            text_box.text_frame.word_wrap = True
            
            # Add image on right side
            image_stream = self._download_image(content.image_url)
            if image_stream:
                img_left = Inches(6.2)
                img_top = Inches(2.0)
                img_width = Inches(3.5)
                
                slide.shapes.add_picture(image_stream, img_left, img_top, width=img_width)
                logger.info("Added text-with-image layout")
        except Exception as e:
            logger.error(f"Failed to create text-with-image layout: {e}")
    
    def _add_list_with_image_layout(self, slide, content, colors):
        """Layout: bullet list on left, image on right."""
        try:
            bg_color = colors.get("background", colors.get("primary", RGBColor(255, 255, 255)))
            auto_text_color = get_contrasting_color(bg_color)
            
            # Add text on left side
            left = Inches(0.8)
            top = Inches(2.0)
            width = Inches(5.0)
            height = Inches(5.0)
            
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            
            # Add bullet points
            for point in content.content:
                p = text_frame.add_paragraph()
                p.text = point
                p.font.size = Pt(18)
                p.font.color.rgb = auto_text_color
                p.space_after = Pt(12)
            
            # Add image on right side
            image_stream = self._download_image(content.image_url)
            if image_stream:
                img_left = Inches(6.2)
                img_top = Inches(2.0)
                img_width = Inches(3.5)
                
                slide.shapes.add_picture(image_stream, img_left, img_top, width=img_width)
                logger.info("Added list-with-image layout")
        except Exception as e:
            logger.error(f"Failed to create list-with-image layout: {e}")
