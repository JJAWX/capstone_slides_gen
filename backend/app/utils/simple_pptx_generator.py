"""
简化的PPTX生成器 - 直接从JSON内容生成PPTX
支持多样化布局：双栏、引用、时间线、对比等
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
from typing import List, Dict, Any, Optional
import logging
import requests
from io import BytesIO
import re

logger = logging.getLogger(__name__)

# 常量定义
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)  # 16:9
MAX_BULLET_CHARS = 80
MAX_BULLETS = 6
MAX_PARAGRAPH_CHARS = 600


class SimplePPTXGenerator:
    """直接从JSON内容生成PPTX，支持多样化布局"""
    
    def __init__(self):
        self.color_schemes = {
            "corporate": {
                "primary": RGBColor(0, 51, 102),
                "secondary": RGBColor(0, 102, 204),
                "accent": RGBColor(255, 153, 0),
                "text": RGBColor(51, 51, 51),
                "light": RGBColor(240, 240, 240),
                "white": RGBColor(255, 255, 255)
            },
            "academic": {
                "primary": RGBColor(51, 51, 51),
                "secondary": RGBColor(102, 102, 102),
                "accent": RGBColor(0, 102, 204),
                "text": RGBColor(0, 0, 0),
                "light": RGBColor(245, 245, 245),
                "white": RGBColor(255, 255, 255)
            },
            "startup": {
                "primary": RGBColor(106, 27, 154),
                "secondary": RGBColor(156, 39, 176),
                "accent": RGBColor(255, 193, 7),
                "text": RGBColor(33, 33, 33),
                "light": RGBColor(243, 229, 245),
                "white": RGBColor(255, 255, 255)
            },
            "minimal": {
                "primary": RGBColor(0, 0, 0),
                "secondary": RGBColor(100, 100, 100),
                "accent": RGBColor(255, 0, 80),
                "text": RGBColor(0, 0, 0),
                "light": RGBColor(250, 250, 250),
                "white": RGBColor(255, 255, 255)
            },
            "creative": {
                "primary": RGBColor(255, 45, 85),
                "secondary": RGBColor(255, 204, 0),
                "accent": RGBColor(88, 86, 214),
                "text": RGBColor(33, 33, 33),
                "light": RGBColor(255, 240, 245),
                "white": RGBColor(255, 255, 255)
            },
            "dark": {
                "primary": RGBColor(138, 180, 248),
                "secondary": RGBColor(255, 121, 198),
                "accent": RGBColor(189, 147, 249),
                "text": RGBColor(205, 214, 244),
                "light": RGBColor(49, 50, 68),
                "white": RGBColor(30, 30, 46)
            }
        }
    
    def generate_pptx(
        self,
        slides_data: List[Dict[str, Any]],
        design_config: Dict[str, Any],
        output_path: Path,
        title: str,
        template: str = "corporate"
    ) -> Path:
        """从幻灯片数据直接生成PPTX"""
        
        logger.info(f"开始生成PPTX: {title}")
        logger.info(f"  - 幻灯片数量: {len(slides_data)}")
        logger.info(f"  - 模板: {template}")
        
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
        
        colors = self._get_colors(design_config, template)
        
        for idx, slide_data in enumerate(slides_data, 1):
            try:
                self._create_slide(prs, slide_data, colors, idx, template)
                if idx % 5 == 0:
                    logger.info(f"  已创建 {idx}/{len(slides_data)} 张幻灯片")
            except Exception as e:
                logger.error(f"  创建幻灯片 {idx} 失败: {e}")
                self._create_error_slide(prs, idx, str(e), colors)
        
        prs.save(str(output_path))
        logger.info(f"PPTX生成完成: {output_path}")
        return output_path
    
    def _get_colors(self, design_config: Dict[str, Any], template: str) -> Dict[str, RGBColor]:
        """获取颜色方案"""
        if design_config and "colors" in design_config:
            colors_data = design_config["colors"]
            try:
                return {
                    "primary": RGBColor(*colors_data.get("primary", [0, 51, 102])),
                    "secondary": RGBColor(*colors_data.get("secondary", [0, 102, 204])),
                    "accent": RGBColor(*colors_data.get("accent", [255, 153, 0])),
                    "text": RGBColor(*colors_data.get("text_main", [51, 51, 51])),
                    "light": RGBColor(240, 240, 240),
                    "white": RGBColor(255, 255, 255)
                }
            except Exception as e:
                logger.warning(f"无法解析design_config颜色: {e}")
        return self.color_schemes.get(template, self.color_schemes["corporate"])
    
    def _truncate_text(self, text: str, max_chars: int) -> str:
        """截断文本"""
        if not text:
            return ""
        if len(text) <= max_chars:
            return text
        return text[:max_chars-3] + "..."
    
    def _get_font_size_for_content(self, content: List[str], base_size: int = 18) -> int:
        """根据内容长度动态调整字体大小"""
        if not content:
            return base_size
        max_len = max(len(c) for c in content)
        total_len = sum(len(c) for c in content)
        if total_len > 400 or max_len > 80:
            return max(12, base_size - 4)
        elif total_len > 300 or max_len > 60:
            return max(14, base_size - 2)
        return base_size
    
    def _create_slide(self, prs: Presentation, slide_data: Dict[str, Any], 
                      colors: Dict[str, RGBColor], slide_number: int, template: str):
        """创建单张幻灯片，根据layout_type选择布局"""
        
        slide_type = slide_data.get("slideType", "content")
        layout_type = slide_data.get("layout_type", "bullet_points")
        
        if slide_type == "title":
            self._create_title_slide(prs, slide_data, colors, template)
        elif layout_type == "section_divider":
            self._create_section_slide(prs, slide_data, colors)
        elif layout_type == "two_column":
            self._create_two_column_slide(prs, slide_data, colors)
        elif layout_type == "comparison":
            self._create_comparison_slide(prs, slide_data, colors)
        elif layout_type == "quote":
            self._create_quote_slide(prs, slide_data, colors)
        elif layout_type == "timeline":
            self._create_timeline_slide(prs, slide_data, colors)
        elif layout_type == "chart_data" and slide_data.get("chart_url"):
            self._create_chart_slide(prs, slide_data, colors)
        elif layout_type == "table_data" or slide_data.get("table"):
            self._create_table_slide(prs, slide_data, colors)
        elif layout_type == "image_content" or slide_data.get("image_url"):
            self._create_image_text_slide(prs, slide_data, colors)
        elif layout_type == "narrative" or slide_data.get("paragraph"):
            self._create_narrative_slide(prs, slide_data, colors)
        else:
            self._create_bullet_slide(prs, slide_data, colors)
    
    def _create_title_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                            colors: Dict[str, RGBColor], template: str):
        """创建主标题页 - 唯一允许背景图的页面"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 只有标题页才添加背景图
        bg_url = slide_data.get("background_image_url")
        if bg_url:
            self._add_background_image(slide, bg_url)
        
        # 半透明遮罩
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.fore_color.brightness = 0.3
        overlay.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = colors["white"]
        p.alignment = PP_ALIGN.CENTER
        
        # 副标题
        content = slide_data.get("content", [])
        if content:
            subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(1))
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = " | ".join(content[:3])
            p.font.size = Pt(20)
            p.font.color.rgb = colors["white"]
            p.alignment = PP_ALIGN.CENTER
    
    def _create_section_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                              colors: Dict[str, RGBColor]):
        """创建小节分隔页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 背景色块
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = colors["primary"]
        bg_shape.line.fill.background()
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.2))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = colors["white"]
        p.alignment = PP_ALIGN.LEFT
        
        # 装饰线
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.3), Inches(2), Inches(0.05))
        line.fill.solid()
        line.fill.fore_color.rgb = colors["accent"]
        line.line.fill.background()
        
        # 要点
        content = slide_data.get("content", [])
        if content:
            points_box = slide.shapes.add_textbox(Inches(1), Inches(3.6), Inches(8), Inches(1.5))
            tf = points_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(content[:4]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = "• " + self._truncate_text(point, 60)
                p.font.size = Pt(18)
                p.font.color.rgb = colors["white"]
                p.space_before = Pt(6)
    
    def _create_bullet_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                             colors: Dict[str, RGBColor]):
        """创建标准bullet points页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        content = slide_data.get("content", [])
        if content:
            font_size = self._get_font_size_for_content(content)
            body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8))
            tf = body_box.text_frame
            tf.word_wrap = True
            
            for i, point in enumerate(content[:MAX_BULLETS]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                text = self._truncate_text(point, MAX_BULLET_CHARS)
                p.text = "• " + text
                p.font.size = Pt(font_size)
                p.font.color.rgb = colors["text"]
                p.space_before = Pt(8)
                p.space_after = Pt(4)
    
    def _create_two_column_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                                  colors: Dict[str, RGBColor]):
        """创建双栏布局页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        content = slide_data.get("content", [])
        left_items = []
        right_items = []
        
        for item in content:
            if item.startswith("LEFT:"):
                left_items.append(item.replace("LEFT:", "").strip())
            elif item.startswith("RIGHT:"):
                right_items.append(item.replace("RIGHT:", "").strip())
            else:
                if len(left_items) <= len(right_items):
                    left_items.append(item)
                else:
                    right_items.append(item)
        
        if not left_items and not right_items and content:
            mid = len(content) // 2
            left_items = content[:mid]
            right_items = content[mid:]
        
        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(3.8))
        self._fill_column(left_box, left_items, colors)
        
        # 分隔线
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(4.95), Inches(1.3), Inches(0.02), Inches(3.5)
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = colors["light"]
        divider.line.fill.background()
        
        # 右栏
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.8))
        self._fill_column(right_box, right_items, colors)
    
    def _fill_column(self, textbox, items: List[str], colors: Dict[str, RGBColor]):
        """填充栏目内容"""
        tf = textbox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items[:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            text = self._truncate_text(item, 50)
            p.text = "• " + text
            p.font.size = Pt(14)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(6)
    
    def _create_comparison_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                                  colors: Dict[str, RGBColor]):
        """创建对比页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        # 左侧卡片
        left_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.4), Inches(4.3), Inches(3.7)
        )
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = colors["light"]
        left_card.line.color.rgb = colors["primary"]
        
        # 右侧卡片
        right_card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.4), Inches(4.3), Inches(3.7)
        )
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = colors["light"]
        right_card.line.color.rgb = colors["secondary"]
        
        content = slide_data.get("content", [])
        mid = len(content) // 2
        
        # 左侧内容
        left_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(3.9), Inches(3.3))
        tf = left_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "方案 A"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
        
        for item in content[:mid][:4]:
            p = tf.add_paragraph()
            p.text = "✓ " + self._truncate_text(item, 40)
            p.font.size = Pt(13)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(6)
        
        # 右侧内容
        right_text = slide.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(3.9), Inches(3.3))
        tf = right_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "方案 B"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors["secondary"]
        
        for item in content[mid:][:4]:
            p = tf.add_paragraph()
            p.text = "✓ " + self._truncate_text(item, 40)
            p.font.size = Pt(13)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(6)
        
        # VS标识
        vs_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.6), Inches(2.7), Inches(0.7), Inches(0.7))
        vs_shape.fill.solid()
        vs_shape.fill.fore_color.rgb = colors["accent"]
        vs_shape.line.fill.background()
        
        vs_text = slide.shapes.add_textbox(Inches(4.6), Inches(2.85), Inches(0.7), Inches(0.4))
        tf = vs_text.text_frame
        p = tf.paragraphs[0]
        p.text = "VS"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = colors["white"]
        p.alignment = PP_ALIGN.CENTER
    
    def _create_quote_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                            colors: Dict[str, RGBColor]):
        """创建引用页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 大引号
        quote_mark = slide.shapes.add_textbox(Inches(0.8), Inches(1), Inches(1), Inches(1))
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(120)
        p.font.color.rgb = colors["accent"]
        
        # 引用文本
        quote_text = slide_data.get("paragraph", "")
        if not quote_text:
            content = slide_data.get("content", [])
            quote_text = content[0] if content else slide_data.get("title", "")
        
        quote_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(2.5))
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self._truncate_text(quote_text, 250)
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = colors["primary"]
        p.alignment = PP_ALIGN.LEFT
        
        # 作者
        author = ""
        content = slide_data.get("content", [])
        if len(content) > 1:
            author = content[-1]
        
        if author:
            author_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.3), Inches(7), Inches(0.5))
            tf = author_box.text_frame
            p = tf.paragraphs[0]
            p.text = "— " + author
            p.font.size = Pt(18)
            p.font.color.rgb = colors["secondary"]
    
    def _create_timeline_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                               colors: Dict[str, RGBColor]):
        """创建时间线页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        content = slide_data.get("content", [])
        num_items = min(len(content), 5)
        
        if num_items == 0:
            return
        
        # 横向时间线
        timeline_y = Inches(3)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.5), timeline_y, Inches(9), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = colors["secondary"]
        line.line.fill.background()
        
        spacing = 8.5 / max(num_items - 1, 1)
        
        for i, item in enumerate(content[:num_items]):
            x = Inches(0.75 + i * spacing)
            
            # 圆点
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x - Inches(0.1), timeline_y - Inches(0.1), Inches(0.25), Inches(0.25)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = colors["accent"]
            dot.line.fill.background()
            
            parts = item.split(":", 1) if ":" in item else [f"Step {i+1}", item]
            date = parts[0].strip()
            desc = parts[1].strip() if len(parts) > 1 else ""
            
            # 日期标签
            date_box = slide.shapes.add_textbox(x - Inches(0.5), timeline_y - Inches(0.7), Inches(1.2), Inches(0.5))
            tf = date_box.text_frame
            p = tf.paragraphs[0]
            p.text = date
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = colors["primary"]
            p.alignment = PP_ALIGN.CENTER
            
            # 描述
            if desc:
                desc_box = slide.shapes.add_textbox(x - Inches(0.6), timeline_y + Inches(0.3), Inches(1.4), Inches(1.2))
                tf = desc_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = self._truncate_text(desc, 40)
                p.font.size = Pt(10)
                p.font.color.rgb = colors["text"]
                p.alignment = PP_ALIGN.CENTER
    
    def _create_narrative_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                                colors: Dict[str, RGBColor]):
        """创建叙述/段落页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        paragraph_text = slide_data.get("paragraph", "")
        if not paragraph_text:
            content = slide_data.get("content", [])
            paragraph_text = " ".join(content)
        
        paragraph_text = self._truncate_text(paragraph_text, MAX_PARAGRAPH_CHARS)
        
        font_size = 16
        if len(paragraph_text) > 400:
            font_size = 14
        elif len(paragraph_text) > 250:
            font_size = 15
        
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(8.4), Inches(3.8))
        tf = body_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = paragraph_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = colors["text"]
        p.line_spacing = 1.5
    
    def _create_image_text_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                                  colors: Dict[str, RGBColor]):
        """创建图文混排页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        image_url = slide_data.get("image_url", "")
        
        if image_url:
            try:
                if image_url.startswith("http"):
                    response = requests.get(image_url, timeout=10)
                    if response.status_code == 200:
                        image_bytes = BytesIO(response.content)
                        slide.shapes.add_picture(image_bytes, Inches(0.5), Inches(1.3), width=Inches(4.5), height=Inches(3.5))
                elif Path(image_url).exists():
                    slide.shapes.add_picture(image_url, Inches(0.5), Inches(1.3), width=Inches(4.5), height=Inches(3.5))
            except Exception as e:
                logger.warning(f"图片加载失败: {e}")
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(4.5), Inches(3.5)
                )
                placeholder.fill.solid()
                placeholder.fill.fore_color.rgb = colors["light"]
        
        content = slide_data.get("content", [])
        text_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.8))
        tf = text_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(content[:5]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + self._truncate_text(point, 50)
            p.font.size = Pt(14)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(8)
    
    def _create_table_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                            colors: Dict[str, RGBColor]):
        """创建表格页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        table_data = slide_data.get("table", {})
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        
        if not headers or not rows:
            content = slide_data.get("content", [])
            if content:
                self._add_bullet_content(slide, content, colors)
            return
        
        headers = headers[:5]
        rows = rows[:6]
        
        rows_count = len(rows) + 1
        cols_count = len(headers)
        
        table_shape = slide.shapes.add_table(
            rows_count, cols_count, Inches(0.5), Inches(1.4), Inches(9), Inches(min(3.5, 0.5 * rows_count))
        )
        table = table_shape.table
        
        for col_idx, header in enumerate(headers):
            cell = table.rows[0].cells[col_idx]
            cell.text = str(header)[:20]
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors["primary"]
            tf = cell.text_frame
            tf.paragraphs[0].font.color.rgb = colors["white"]
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        for row_idx, row_data in enumerate(rows, 1):
            for col_idx, cell_data in enumerate(row_data[:cols_count]):
                cell = table.rows[row_idx].cells[col_idx]
                cell.text = str(cell_data)[:25]
                tf = cell.text_frame
                tf.paragraphs[0].font.size = Pt(11)
                tf.paragraphs[0].font.color.rgb = colors["text"]
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    def _create_chart_slide(self, prs: Presentation, slide_data: Dict[str, Any],
                            colors: Dict[str, RGBColor]):
        """创建图表页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, slide_data.get("title", ""), colors)
        
        chart_url = slide_data.get("chart_url", "")
        if chart_url and Path(chart_url).exists():
            try:
                slide.shapes.add_picture(chart_url, Inches(1), Inches(1.2), width=Inches(8))
            except Exception as e:
                logger.error(f"图表添加失败: {e}")
        
        content = slide_data.get("content", [])
        if content:
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.8))
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = " | ".join(content[:3])
            p.font.size = Pt(12)
            p.font.color.rgb = colors["secondary"]
            p.alignment = PP_ALIGN.CENTER
    
    def _add_slide_title(self, slide, title: str, colors: Dict[str, RGBColor]):
        """添加幻灯片标题"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = self._truncate_text(title, 60)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = colors["primary"]
    
    def _add_bullet_content(self, slide, content: List[str], colors: Dict[str, RGBColor]):
        """添加bullet内容"""
        body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(3.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(content[:MAX_BULLETS]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + self._truncate_text(point, MAX_BULLET_CHARS)
            p.font.size = Pt(16)
            p.font.color.rgb = colors["text"]
            p.space_before = Pt(6)
    
    def _add_background_image(self, slide, bg_url: str):
        """添加背景图片 - 仅用于标题页"""
        try:
            response = requests.get(bg_url, timeout=10)
            if response.status_code == 200:
                image_bytes = BytesIO(response.content)
                pic = slide.shapes.add_picture(image_bytes, Inches(0), Inches(0), width=SLIDE_WIDTH, height=SLIDE_HEIGHT)
                spTree = slide.shapes._spTree
                spTree.remove(pic._element)
                spTree.insert(2, pic._element)
                logger.info("  添加背景图片成功")
        except Exception as e:
            logger.warning(f"  背景图片加载失败: {e}")
    
    def _create_error_slide(self, prs: Presentation, slide_number: int, error_message: str,
                            colors: Dict[str, RGBColor]):
        """创建错误占位幻灯片"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Slide {slide_number} - Generation Error"
        p.font.size = Pt(24)
        p.font.color.rgb = colors["accent"]
        p.alignment = PP_ALIGN.CENTER
