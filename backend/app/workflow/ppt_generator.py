import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from typing import Dict, List, Tuple
import os

class XMLTemplateReader:
    def __init__(self, xml_file_path: str):
        self.xml_file_path = xml_file_path
        self.placeholders = []
        self._parse_xml()

    def _parse_xml(self):
        tree = ET.parse(self.xml_file_path)
        root = tree.getroot()
        
        # Find all shape elements that have placeholders
        for sp in root.findall('.//p:sp', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}):
            placeholder_info = self._extract_placeholder_info(sp)
            if placeholder_info:
                self.placeholders.append(placeholder_info)

    def _extract_placeholder_info(self, sp_element) -> Dict:
        # Check if this shape has a placeholder
        nvpr = sp_element.find('.//p:nvPr', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if nvpr is None:
            return None
            
        ph = nvpr.find('p:ph', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if ph is None:
            return None
            
        ph_type = ph.get('type')
        idx = ph.get('idx')
        
        # Find transform element for position/size
        sppr = sp_element.find('p:spPr', {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'})
        if sppr is not None:
            xfrm = sppr.find('a:xfrm', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
            if xfrm is not None:
                off = xfrm.find('a:off', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                ext = xfrm.find('a:ext', {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                
                if off is not None and ext is not None:
                    x = int(off.get('x', 0))
                    y = int(off.get('y', 0))
                    cx = int(ext.get('cx', 0))
                    cy = int(ext.get('cy', 0))
                    
                    # Convert EMU to inches
                    x_in = x / 914400
                    y_in = y / 914400
                    cx_in = cx / 914400
                    cy_in = cy / 914400
                    
                    return {
                        'type': ph_type,
                        'idx': idx,
                        'x': x_in,
                        'y': y_in,
                        'cx': cx_in,
                        'cy': cy_in
                    }
        return None

class PPTGenerator:
    def __init__(self, template_reader: XMLTemplateReader):
        self.template_reader = template_reader
        self.presentation = Presentation()
        # 设置为16:9尺寸 (10英寸 x 5.625英寸)
        self.presentation.slide_width = Inches(10)
        self.presentation.slide_height = Inches(5.625)

    def generate_slide(self, content: Dict[str, str]) -> None:
        slide_layout = self.presentation.slide_layouts[5]  # Blank layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        for placeholder in self.template_reader.placeholders:
            self._fill_placeholder(slide, placeholder, content)

    def _fill_placeholder(self, slide, placeholder: Dict, content: Dict[str, str]):
        ph_type = placeholder['type']
        idx = placeholder['idx']
        
        if ph_type == 'title':
            title = slide.shapes.title
            if title:
                title.text = content.get('title', '')
                title.left = Inches(placeholder['x'])
                title.top = Inches(placeholder['y'])
                title.width = Inches(placeholder['cx'])
                title.height = Inches(placeholder['cy'])
        elif ph_type == 'body':
            # For body placeholders, add text box
            left = Inches(placeholder['x'])
            top = Inches(placeholder['y'])
            width = Inches(placeholder['cx'])
            height = Inches(placeholder['cy'])
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = content.get(f'body_{idx}', '')
        elif ph_type == 'pic':
            # For picture placeholders, add picture with aspect ratio preservation
            image_path = content.get(f'image_{idx}')
            if image_path and os.path.exists(image_path):
                self._add_picture_with_aspect_ratio(slide, image_path, placeholder)

    def _add_picture_with_aspect_ratio(self, slide, image_path: str, placeholder: Dict):
        """Add picture while preserving aspect ratio and centering it in the placeholder area"""
        try:
            from PIL import Image
            
            # Get placeholder dimensions
            ph_left = Inches(placeholder['x'])
            ph_top = Inches(placeholder['y'])
            ph_width = Inches(placeholder['cx'])
            ph_height = Inches(placeholder['cy'])
            
            # Get image dimensions
            with Image.open(image_path) as img:
                img_width, img_height = img.size
            
            # Calculate aspect ratios
            img_aspect = img_width / img_height
            ph_aspect = ph_width / ph_height
            
            # Calculate new dimensions to fit within placeholder while preserving aspect ratio
            if img_aspect > ph_aspect:
                # Image is wider relative to placeholder - fit by width
                new_width = ph_width
                new_height = ph_width / img_aspect
            else:
                # Image is taller relative to placeholder - fit by height
                new_height = ph_height
                new_width = ph_height * img_aspect
            
            # Center the image in the placeholder
            new_left = ph_left + (ph_width - new_width) / 2
            new_top = ph_top + (ph_height - new_height) / 2
            
            # Add the picture
            slide.shapes.add_picture(image_path, new_left, new_top, new_width, new_height)
            
        except ImportError:
            # Fallback if PIL is not available - use original method
            left = Inches(placeholder['x'])
            top = Inches(placeholder['y'])
            width = Inches(placeholder['cx'])
            height = Inches(placeholder['cy'])
            slide.shapes.add_picture(image_path, left, top, width, height)
        except Exception:
            # If anything goes wrong, skip adding the picture
            pass

    def save(self, output_path: str):
        self.presentation.save(output_path)

def generate_pptx_from_xml(xml_file_path: str, content: Dict[str, str], output_dir: str = 'outputs/pptx') -> str:
    """
    Generate PPTX from XML template
    
    Args:
        xml_file_path: Path to XML template file
        content: Dictionary with content to fill placeholders
        output_dir: Output directory for PPTX files
    
    Returns:
        Path to generated PPTX file
    """
    # Ensure output directory exists
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Parse XML template
    reader = XMLTemplateReader(xml_file_path)
    
    # Generate PPT
    generator = PPTGenerator(reader)
    generator.generate_slide(content)
    
    # Generate output filename
    template_name = Path(xml_file_path).stem
    output_path = Path(output_dir) / f"{template_name}.pptx"
    
    # Save PPTX
    generator.save(str(output_path))
    
    return str(output_path)

if __name__ == "__main__":
    # Example usage
    xml_path = "templates/content_slide.xml"
    content = {
        'title': 'Sample Title',
        'body_1': 'Sample content for body placeholder'
    }
    output = generate_pptx_from_xml(xml_path, content)
    print(f"Generated PPTX: {output}")